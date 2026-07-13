"""PPTX backend: transpiles OVP reference SVGs into native PowerPoint shapes.

One converter covers every chart, component, and recipe, current and
future, because all OVP renders use the same drawing vocabulary:
rect, line, polyline, polygon, circle, text, nested svg viewport.

Determinism contract for PPTX: the zip container carries volatile
timestamps, so the conformance check compares every MEMBER of the
package (slide XML, relationships, content types) instead of raw bytes.
Same SVG -> same members, always.

    python tools/render_pptx.py --svg golden/RC-001/RC-001.svg \
        --out golden/RC-001/RC-001.pptx
"""

import argparse
import io
import re
import xml.etree.ElementTree as ET
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

EMU_PER_PX = 9525  # 96 dpi
NS = "{http://www.w3.org/2000/svg}"


def hex_rgb(h):
    from pptx.dml.color import RGBColor
    return RGBColor(int(h[1:3], 16), int(h[3:5], 16), int(h[5:7], 16))


def first_family(stack):
    fam = stack.split(",")[0].strip()
    return fam.strip("'\"")


def set_fill_alpha(shape, alpha_pct_100k):
    sf = shape.fill._xPr.find(qn("a:solidFill"))
    clr = sf.find(qn("a:srgbClr"))
    a = clr.makeelement(qn("a:alpha"), {"val": str(alpha_pct_100k)})
    clr.append(a)


def no_line(shape):
    shape.line.fill.background()


def solid(shape, hexv):
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(hexv)


class T:
    """Affine transform (uniform scale + translate) for nested viewports."""

    def __init__(self, s=1.0, tx=0.0, ty=0.0):
        self.s, self.tx, self.ty = s, tx, ty

    def x(self, v):
        return float(v) * self.s + self.tx

    def y(self, v):
        return float(v) * self.s + self.ty

    def d(self, v):
        return float(v) * self.s


def add_rect(shapes, el, t):
    x, y = t.x(el.get("x", 0)), t.y(el.get("y", 0))
    w, h = t.d(el.get("width")), t.d(el.get("height"))
    rx = float(el.get("rx", 0))
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if rx else MSO_SHAPE.RECTANGLE
    sp = shapes.add_shape(kind, Emu(int(x * EMU_PER_PX)), Emu(int(y * EMU_PER_PX)),
                          Emu(int(w * EMU_PER_PX)), Emu(int(h * EMU_PER_PX)))
    if rx:
        sp.adjustments[0] = min(0.5, t.d(rx) / min(w, h))
    fill = el.get("fill", "none")
    if fill != "none":
        solid(sp, fill)
        if el.get("fill-opacity"):
            set_fill_alpha(sp, int(float(el.get("fill-opacity")) * 100000))
    else:
        sp.fill.background()
    stroke = el.get("stroke")
    if stroke and stroke != "none":
        sp.line.color.rgb = hex_rgb(stroke)
        sp.line.width = Pt(float(el.get("stroke-width", 1)) * t.s * 0.75)
    else:
        no_line(sp)
    sp.shadow.inherit = False


def add_circle(shapes, el, t):
    cx, cy, r = t.x(el.get("cx")), t.y(el.get("cy")), t.d(el.get("r"))
    sp = shapes.add_shape(MSO_SHAPE.OVAL,
                          Emu(int((cx - r) * EMU_PER_PX)), Emu(int((cy - r) * EMU_PER_PX)),
                          Emu(int(2 * r * EMU_PER_PX)), Emu(int(2 * r * EMU_PER_PX)))
    solid(sp, el.get("fill"))
    no_line(sp)
    sp.shadow.inherit = False


def add_line(shapes, el, t):
    x1, y1 = t.x(el.get("x1")), t.y(el.get("y1"))
    x2, y2 = t.x(el.get("x2")), t.y(el.get("y2"))
    conn = shapes.add_connector(1, Emu(int(x1 * EMU_PER_PX)), Emu(int(y1 * EMU_PER_PX)),
                                Emu(int(x2 * EMU_PER_PX)), Emu(int(y2 * EMU_PER_PX)))
    conn.line.color.rgb = hex_rgb(el.get("stroke"))
    conn.line.width = Pt(float(el.get("stroke-width", 1)) * t.s * 0.75)
    if el.get("stroke-dasharray"):
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    conn.shadow.inherit = False


def add_poly(shapes, el, t, close):
    pts = [(t.x(a), t.y(b)) for a, b in
           (p.split(",") for p in el.get("points").split())]
    fb = shapes.build_freeform(pts[0][0], pts[0][1], scale=EMU_PER_PX)
    fb.add_line_segments(pts[1:], close=close)
    sp = fb.convert_to_shape()
    fill = el.get("fill", "none")
    if close and fill != "none":
        solid(sp, fill)
        if el.get("fill-opacity"):
            set_fill_alpha(sp, int(float(el.get("fill-opacity")) * 100000))
    else:
        sp.fill.background()
    stroke = el.get("stroke")
    if stroke and stroke != "none":
        sp.line.color.rgb = hex_rgb(stroke)
        sp.line.width = Pt(float(el.get("stroke-width", 1)) * t.s * 0.75)
    else:
        no_line(sp)
    sp.shadow.inherit = False


def _flatten_path(d):
    """SVG path -> list of subpaths (point lists). Arcs and beziers are
    flattened deterministically (fixed sample counts), covering the OVP
    drawing vocabulary: M, L, A (circular), Q, C, Z."""
    import math
    tokens = d.replace(",", " ").split()
    i = 0
    subpaths = []
    pts = []
    cur = (0.0, 0.0)

    def num():
        nonlocal i
        v = float(tokens[i])
        i += 1
        return v

    while i < len(tokens):
        cmd = tokens[i]
        i += 1
        if cmd == "M":
            if pts:
                subpaths.append((pts, False))
            cur = (num(), num())
            pts = [cur]
        elif cmd == "L":
            cur = (num(), num())
            pts.append(cur)
        elif cmd == "A":
            rx = num(); num()  # ry == rx in OVP output
            num()  # rotation, always 0
            large = int(num()); sweep = int(num())
            end = (num(), num())
            x1, y1 = cur; x2, y2 = end
            mx, my = (x1 + x2) / 2, (y1 + y2) / 2
            dx, dy = x2 - x1, y2 - y1
            q = math.hypot(dx, dy)
            if q == 0 or rx <= 0:
                pts.append(end); cur = end; continue
            h2 = max(rx * rx - (q / 2) ** 2, 0.0)
            h = math.sqrt(h2)
            ux, uy = -dy / q, dx / q
            if (sweep == 1) != (large == 1):
                cxp, cyp = mx + ux * h, my + uy * h
            else:
                cxp, cyp = mx - ux * h, my - uy * h
            a1 = math.atan2(y1 - cyp, x1 - cxp)
            a2 = math.atan2(y2 - cyp, x2 - cxp)
            if sweep == 1:
                while a2 < a1:
                    a2 += 2 * math.pi
            else:
                while a2 > a1:
                    a2 -= 2 * math.pi
            steps = 24
            for s in range(1, steps + 1):
                a = a1 + (a2 - a1) * s / steps
                pts.append((cxp + rx * math.cos(a), cyp + rx * math.sin(a)))
            cur = end
        elif cmd == "Q":
            c1 = (num(), num()); end = (num(), num())
            x0, y0 = cur
            for s in range(1, 17):
                u = s / 16
                x = (1 - u) ** 2 * x0 + 2 * (1 - u) * u * c1[0] + u * u * end[0]
                y = (1 - u) ** 2 * y0 + 2 * (1 - u) * u * c1[1] + u * u * end[1]
                pts.append((x, y))
            cur = end
        elif cmd == "C":
            c1 = (num(), num()); c2 = (num(), num()); end = (num(), num())
            x0, y0 = cur
            for s in range(1, 17):
                u = s / 16
                x = ((1 - u) ** 3 * x0 + 3 * (1 - u) ** 2 * u * c1[0]
                     + 3 * (1 - u) * u * u * c2[0] + u ** 3 * end[0])
                y = ((1 - u) ** 3 * y0 + 3 * (1 - u) ** 2 * u * c1[1]
                     + 3 * (1 - u) * u * u * c2[1] + u ** 3 * end[1])
                pts.append((x, y))
            cur = end
        elif cmd == "Z":
            subpaths.append((pts, True))
            pts = []
    if pts:
        subpaths.append((pts, False))
    return subpaths


def add_path(shapes, el, t):
    for pts, closed in _flatten_path(el.get("d")):
        if len(pts) < 2:
            continue
        tp = [(t.x(a), t.y(b)) for a, b in pts]
        fb = shapes.build_freeform(tp[0][0], tp[0][1], scale=EMU_PER_PX)
        fb.add_line_segments(tp[1:], close=closed)
        sp = fb.convert_to_shape()
        fill = el.get("fill", "none")
        if closed and fill != "none":
            solid(sp, fill)
            if el.get("fill-opacity"):
                set_fill_alpha(sp, int(float(el.get("fill-opacity")) * 100000))
        else:
            sp.fill.background()
        stroke = el.get("stroke")
        if stroke and stroke != "none":
            sp.line.color.rgb = hex_rgb(stroke)
            sp.line.width = Pt(float(el.get("stroke-width", 1)) * t.s * 0.75)
            if el.get("stroke-opacity"):
                pass  # line alpha not part of the conformance contract
            if el.get("stroke-dasharray"):
                from pptx.enum.dml import MSO_LINE_DASH_STYLE
                sp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        else:
            no_line(sp)
        sp.shadow.inherit = False


def add_text(shapes, el, t):
    size_px = float(el.get("font-size")) * t.s
    x, y = t.x(el.get("x")), t.y(el.get("y"))
    anchor = el.get("text-anchor", "start")
    box_w = 600.0 * t.s
    if el.get("dominant-baseline") == "middle":
        top = y - size_px * 0.65
    else:
        top = y - size_px
    if anchor == "start":
        left, align = x, PP_ALIGN.LEFT
    elif anchor == "middle":
        left, align = x - box_w / 2, PP_ALIGN.CENTER
    else:
        left, align = x - box_w, PP_ALIGN.RIGHT
    tb = shapes.add_textbox(Emu(int(left * EMU_PER_PX)), Emu(int(top * EMU_PER_PX)),
                            Emu(int(box_w * EMU_PER_PX)),
                            Emu(int(size_px * 1.5 * EMU_PER_PX)))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = el.text or ""
    f = run.font
    f.size = Pt(round(size_px * 0.75, 1))
    f.name = first_family(el.get("font-family", "Segoe UI"))
    f.color.rgb = hex_rgb(el.get("fill"))
    f.bold = int(el.get("font-weight", 400)) >= 600
    if el.get("font-style") == "italic":
        f.italic = True


def walk(el, shapes, t):
    for child in el:
        tag = child.tag.replace(NS, "")
        if tag == "rect":
            add_rect(shapes, child, t)
        elif tag == "line":
            add_line(shapes, child, t)
        elif tag == "polyline":
            add_poly(shapes, child, t, close=False)
        elif tag == "polygon":
            add_poly(shapes, child, t, close=True)
        elif tag == "circle":
            add_circle(shapes, child, t)
        elif tag == "path":
            add_path(shapes, child, t)
        elif tag == "text":
            add_text(shapes, child, t)
        elif tag == "svg":
            vb = [float(v) for v in child.get("viewBox").split()]
            s = float(child.get("width")) / vb[2]
            nt = T(s * t.s,
                   t.x(child.get("x")) - vb[0] * s * t.s,
                   t.y(child.get("y")) - vb[1] * s * t.s)
            walk(child, shapes, nt)


def convert(svg_text):
    """SVG string -> pptx bytes."""
    root = ET.fromstring(svg_text)
    W, H = float(root.get("width")), float(root.get("height"))
    prs = Presentation()
    prs.slide_width = Emu(int(W * EMU_PER_PX))
    prs.slide_height = Emu(int(H * EMU_PER_PX))
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    walk(root, slide.shapes, T())
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--svg", required=True)
    ap.add_argument("--out", required=True)
    a = ap.parse_args()
    data = convert(Path(a.svg).read_text(encoding="utf-8"))
    Path(a.out).write_bytes(data)
    print(f"rendered {a.out} ({len(data)} bytes)")


if __name__ == "__main__":
    main()
